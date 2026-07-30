import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# ─── COLOR PALETTE ────────────
BG_COLOR       = RGBColor(7, 9, 14)       # #07090E
CARD_BG        = RGBColor(19, 27, 46)     # #131B2E
CARD_BORDER    = RGBColor(30, 41, 59)     # #1E293B
NEON_CYAN      = RGBColor(0, 240, 255)    # #00F0FF
EMERALD_GREEN  = RGBColor(0, 255, 170)    # #00FFAA
PURPLE_ACCENT  = RGBColor(155, 89, 255)   # #9B59FF
GOLD_ACCENT    = RGBColor(255, 179, 0)    # #FFB300
RED_ACCENT     = RGBColor(239, 68, 68)    # #EF4444
WHITE          = RGBColor(255, 255, 255)
TEXT_MUTED     = RGBColor(148, 163, 184)  # #94A3B8

def create_deck():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout     = prs.slide_layouts[6] # blank slide layout

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, slide_num, category, title):
        set_background(slide)
        
        # Category badge
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"ECOLABEL X  •  SLIDE {slide_num:02d}  •  {category.upper()}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = NEON_CYAN

        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE

    def set_speaker_notes(slide, say, click, pause, judges):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = (
            f"🎤 WHAT TO SAY:\n{say}\n\n"
            f"🖱️ WHAT TO CLICK:\n{click}\n\n"
            f"⏸️ WHEN TO PAUSE:\n{pause}\n\n"
            f"💡 WHAT JUDGES SHOULD NOTICE:\n{judges}"
        )

    def add_card(slide, left, top, width, height, title="", subtitle="", bg=CARD_BG, border=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = border
        shape.line.width = Pt(1)

        if title or subtitle:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.2)
            tf.margin_right = Inches(0.2)
            
            if title:
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = WHITE
            
            if subtitle:
                p2 = tf.add_paragraph()
                p2.text = subtitle
                p2.font.size = Pt(11)
                p2.font.color.rgb = TEXT_MUTED
                p2.space_before = Pt(4)
        return shape

    def add_animation(slide, shape):
        sld = slide._element
        timing = sld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
        if timing is None:
            timing = parse_xml(
                '<p:timing %s>'
                '  <p:tnLst>'
                '    <p:par>'
                '      <p:cTn id="1" fill="hold" subTnLst="0" dur="indefinite">'
                '        <p:childTnLst/>'
                '      </p:cTn>'
                '    </p:par>'
                '  </p:tnLst>'
                '</p:timing>' % nsdecls('p')
            )
            sld.append(timing)
        
        child_tn_lst = timing.xpath('.//p:childTnLst')
        if child_tn_lst:
            shp_id = shape.shape_id
            click_xml = parse_xml(
                '<p:seq %s concurrent="1" nextAc="seek">'
                '  <p:cTn id="%d" restart="whenNotActive" fill="hold" evt="onNext" nodeType="mainSeq">'
                '    <p:childTnLst>'
                '      <p:set>'
                '        <p:cb>'
                '          <p:cBhvr>'
                '            <p:cTn id="%d" dur="1" fill="hold">'
                '              <p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                '            </p:cTn>'
                '            <p:tgtEl><p:spTarget spid="%d"/></p:tgtEl>'
                '            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
                '          </p:cBhvr>'
                '          <p:to><p:strVal value="visible"/></p:to>'
                '        </p:cb>'
                '      </p:set>'
                '    </p:childTnLst>'
                '  </p:cTn>'
                '</p:seq>' % (nsdecls('p'), shp_id + 100, shp_id + 200, shp_id)
            )
            child_tn_lst[0].append(click_xml)

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 1: Cinematic Cover
    # ───────────────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1)

    c1 = add_card(s1, 1.5, 1.2, 10.33, 5.1, bg=RGBColor(15, 23, 42), border=NEON_CYAN)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    tf1.margin_top = Inches(0.8)
    
    p = tf1.paragraphs[0]
    p.text = "EcoLabel X"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = EMERALD_GREEN
    p.alignment = PP_ALIGN.CENTER

    p_sub = tf1.add_paragraph()
    p_sub.text = "AI-Powered Sustainability Intelligence Platform"
    p_sub.font.size = Pt(22)
    p_sub.font.bold = True
    p_sub.font.color.rgb = NEON_CYAN
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.space_before = Pt(12)

    p_desc = tf1.add_paragraph()
    p_desc.text = "Enterprise Anti-Greenwashing Audit & Carbon Accounting Engine"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.alignment = PP_ALIGN.CENTER
    p_desc.space_before = Pt(24)

    p_hack = tf1.add_paragraph()
    p_hack.text = "HACKATHON FINALIST PRESENTATION  •  2026"
    p_hack.font.size = Pt(11)
    p_hack.font.bold = True
    p_hack.font.color.rgb = PURPLE_ACCENT
    p_hack.alignment = PP_ALIGN.CENTER
    p_hack.space_before = Pt(36)

    set_speaker_notes(
        s1,
        say="Good morning judges! Today we present EcoLabel X, an enterprise AI platform that solves greenwashing and automates ESG audits in under 1.5 seconds.",
        click="Click once to advance to Slide 2.",
        pause="Pause 2 seconds after reading the title to let judges take in the branding.",
        judges="Judges should notice the sleek Fortune 500 dark theme and clear product positioning."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 2: Problem Statement
    # ───────────────────────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, 2, "Market Friction", "The Core Problem in Sustainability Disclosures")

    problems = [
        ("Massive Volume", "Companies produce 100+ page ESG PDFs filled with complex tables and unstructured claims."),
        ("Manual Friction", "Compliance officers take weeks to manually cross-reference claims against ISO standards."),
        ("Surge in Greenwashing", "Over 40% of environmental claims online and in disclosures are vague or unsubstantiated."),
        ("Regulatory Mandates", "Strict enforcement from EU ESPR & FTC Green Guides requires quantitative evidence.")
    ]

    for i, (title, desc) in enumerate(problems):
        row, col = divmod(i, 2)
        card = add_card(s2, 0.8 + col * 5.8, 1.8 + row * 2.5, 5.5, 2.2, f"❌ {title}", desc, border=RED_ACCENT)
        add_animation(s2, card)

    set_speaker_notes(
        s2,
        say="Sustainability reporting is broken. Companies produce 100-page disclosures, auditors spend weeks checking claims, and greenwashing is at an all-time high.",
        click="Click 4 times — one click per problem card.",
        pause="Pause after revealing 'Regulatory Mandates' to emphasize urgency.",
        judges="Judges will see that we understand the quantitative friction in corporate compliance."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 3: Solution
    # ───────────────────────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, 3, "Our Innovation", "Introducing EcoLabel X: Instant ESG Intelligence")

    solutions = [
        ("One-Click PDF Upload", "Parses multi-page ESG disclosures in seconds.", EMERALD_GREEN),
        ("AI Claim Extraction", "Identifies exact quotes, page numbers, and keywords.", NEON_CYAN),
        ("4-Agent AI Network", "Verification, Carbon, Compliance, and Risk agents run in parallel.", PURPLE_ACCENT),
        ("Audit-Ready Reports", "Generates publication-ready executive PDF audit reports.", GOLD_ACCENT)
    ]

    for i, (title, desc, color) in enumerate(solutions):
        card = add_card(s3, 0.8 + i * 2.9, 1.8, 2.7, 4.8, f"⚡ {title}", desc, border=color)
        add_animation(s3, card)

    set_speaker_notes(
        s3,
        say="EcoLabel X turns weeks of manual auditing into a 1.5-second automated workflow. Upload a PDF, and our 4 AI Agents verify claims and generate compliance reports.",
        click="Click 4 times to reveal each solution pillar sequentially.",
        pause="Pause after the 4th pillar to let the complete solution vision land.",
        judges="Judges will notice end-to-end automation from PDF input to executive PDF export."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 4: System Workflow
    # ───────────────────────────────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, 4, "Execution Flow", "8-Stage Pipeline Workflow")

    steps = [
        ("1. Upload PDF", "Drag & drop report"),
        ("2. PDF Processing", "PyPDF2 text parsing"),
        ("3. Text Extraction", "Clean structure"),
        ("4. Claim Detection", "Regex & NLP tags"),
        ("5. 4 AI Agents", "Parallel analysis"),
        ("6. Results Engine", "Score calculation"),
        ("7. Dashboard", "Live visual data"),
        ("8. Audit Report", "Executive PDF export")
    ]

    for i, (title, desc) in enumerate(steps):
        row, col = divmod(i, 4)
        card = add_card(s4, 0.8 + col * 2.9, 1.8 + row * 2.5, 2.7, 2.1, title, desc, border=NEON_CYAN)
        add_animation(s4, card)

    set_speaker_notes(
        s4,
        say="Here is our end-to-end execution flow. From raw PDF upload to text extraction, 4-agent parallel evaluation, and instant PDF report generation.",
        click="Click 8 times — one step per click to build the pipeline progressively.",
        pause="Pause after step 5 to explain the 4 AI Agents orchestration.",
        judges="Judges will appreciate the clean, modular architecture build."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 5: Overall Technical Architecture
    # ───────────────────────────────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, 5, "System Design", "High-Performance Technical Architecture")

    arch_nodes = [
        ("Next.js 14 Web App", "React, TypeScript, Glassmorphism UI", EMERALD_GREEN, 0.8, 1.8),
        ("FastAPI Backend", "Python 3.10, Uvicorn Async Server", NEON_CYAN, 4.8, 1.8),
        ("PyPDF2 Text Engine", "Raw PDF structure & table reader", PURPLE_ACCENT, 8.8, 1.8),
        ("Google Gemini 2.5 Flash", "LLM reasoning + ISO 14021 rules", GOLD_ACCENT, 0.8, 4.4),
        ("Carbon & Risk Calculators", "Scope 1-3 & Trust Score matrix", EMERALD_GREEN, 4.8, 4.4),
        ("Executive PDF Generator", "Publication-ready A4 exporter", NEON_CYAN, 8.8, 4.4),
    ]

    for title, desc, color, left, top in arch_nodes:
        card = add_card(s5, left, top, 3.7, 2.2, title, desc, border=color)
        add_animation(s5, card)

    set_speaker_notes(
        s5,
        say="Our technical architecture connects Next.js 14 on the frontend to FastAPI on the backend, processing PDFs via PyPDF2 and Gemini 2.5 Flash in under 1.5 seconds.",
        click="Click 6 times to reveal each architectural component block.",
        pause="Pause on Gemini 2.5 Flash to highlight LLM reasoning combined with deterministic fallback rules.",
        judges="Judges will observe technical rigor, fast async performance, and enterprise scalability."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 6: Four AI Agents Overview
    # ───────────────────────────────────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, 6, "AI Core Network", "Multi-Agent Orchestration Engine")

    agents = [
        ("Verification Agent", "Cross-references 400+ certification databases & ISO standards.", EMERALD_GREEN),
        ("Carbon Agent", "Scope 1, 2 & 3 lifecycle carbon accounting engine.", NEON_CYAN),
        ("Compliance Agent", "EU ESPR & FTC Green Guide regulatory readiness tracker.", PURPLE_ACCENT),
        ("Risk Agent", "Detects vague, misleading, and unsubstantiated claims.", RED_ACCENT)
    ]

    for i, (title, desc, color) in enumerate(agents):
        row, col = divmod(i, 2)
        card = add_card(s6, 0.8 + col * 5.8, 1.8 + row * 2.5, 5.5, 2.2, f"🤖 {title}", desc, border=color)
        add_animation(s6, card)

    set_speaker_notes(
        s6,
        say="EcoLabel X deploys four specialized AI agents that run concurrently on every uploaded document: Verification, Carbon, Compliance, and Risk.",
        click="Click 4 times — reveal one agent card per click.",
        pause="Pause after revealing all 4 agents to summarize their synergy.",
        judges="Judges will notice multi-agent architecture handling distinct domain responsibilities."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 7: Verification Agent Deep-Dive
    # ───────────────────────────────────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, 7, "Agent Deep-Dive", "Verification Agent: ISO & Database Matching")

    v_details = [
        ("Purpose", "Cross-references disclosures against ISO 14021 & certification registries.", EMERALD_GREEN, 0.8, 1.8),
        ("Workflow", "Identifies claim -> Extracts evidence -> Matches issuer -> Assigns status.", NEON_CYAN, 6.8, 1.8),
        ("Status Metrics", "Verified (Confidence > 75%) | Under Review (50-75%) | Unverified (<50%).", PURPLE_ACCENT, 0.8, 4.4),
        ("Output Example", "ISO 14064 GHG Verified (TÜV SÜD) -> Verified (Page 2).", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in v_details:
        card = add_card(s7, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s7, card)

    set_speaker_notes(
        s7,
        say="The Verification Agent analyzes exact claim quotes, page references, and assigns confidence scores from Verified to Under Review.",
        click="Click 4 times to reveal each section of the Verification Agent deep-dive.",
        pause="Pause on Status Metrics to explain confidence thresholds.",
        judges="Judges will note rigorous evidence-backed verification rather than simple sentiment analysis."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 8: Carbon Agent Deep-Dive
    # ───────────────────────────────────────────────────────────────────────────
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, 8, "Agent Deep-Dive", "Carbon Agent: Scope 1, 2 & 3 Accounting Engine")

    scopes = [
        ("Scope 1: Direct Operations", "Onsite manufacturing, stationary combustion, fleet operations.", EMERALD_GREEN, 0.8, 1.8),
        ("Scope 2: Energy & Electricity", "Purchased electricity, steam, facility heating & cooling.", NEON_CYAN, 4.8, 1.8),
        ("Scope 3: Value Chain", "Raw material procurement, freight distribution, end-of-life.", PURPLE_ACCENT, 8.8, 1.8),
        ("YoY Intensity Reductions", "Calculates percentage reductions across packaging & supply chain.", GOLD_ACCENT, 0.8, 4.4)
    ]

    for i, (title, desc, color, left, top) in enumerate(scopes):
        w = 11.7 if i == 3 else 3.7
        card = add_card(s8, left, top, w, 2.2, title, desc, border=color)
        add_animation(s8, card)

    set_speaker_notes(
        s8,
        say="Our Carbon Agent parses Scope 1 direct, Scope 2 energy, and Scope 3 value chain disclosures, computing total kilotonnes of CO2e automatically.",
        click="Click 4 times — one per scope and the YoY intensity calculator.",
        pause="Pause on Scope 3 to highlight value chain calculation complexity.",
        judges="Judges will appreciate automated GHG Protocol alignment."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 9: Compliance Agent Deep-Dive
    # ───────────────────────────────────────────────────────────────────────────
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, 9, "Agent Deep-Dive", "Compliance Agent: EU ESPR & FTC Green Guides")

    comp_items = [
        ("EU ESPR Monitoring", "Ecodesign for Sustainable Products Regulation audit readiness check.", EMERALD_GREEN, 0.8, 1.8),
        ("FTC Green Guide Alignment", "Ensures environmental claims meet US Federal Trade Commission standards.", NEON_CYAN, 6.8, 1.8),
        ("Audit Compliance Index", "Calculates dynamic compliance score (0-100%) based on verified claims ratio.", PURPLE_ACCENT, 0.8, 4.4),
        ("Regulatory Checklist", "Flags missing mandatory quantitative proof before submission.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in comp_items:
        card = add_card(s9, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s9, card)

    set_speaker_notes(
        s9,
        say="The Compliance Agent monitors EU ESPR and FTC Green Guide requirements, generating an Audit Compliance Index score to prevent legal exposure.",
        click="Click 4 times to reveal each compliance check module.",
        pause="Pause on Audit Compliance Index to explain compliance readiness scoring.",
        judges="Judges will see immediate real-world regulatory utility for enterprise risk teams."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 10: Greenwashing Risk Agent Deep-Dive
    # ───────────────────────────────────────────────────────────────────────────
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, 10, "Agent Deep-Dive", "Risk Agent: Misleading Claim & Vague Term Detector")

    risk_items = [
        ("Vague Term Detection", "Flags buzzwords like 'eco-friendly', 'green', and 'pure' without data.", RED_ACCENT, 0.8, 1.8),
        ("Missing Baseline Flags", "Detects reduction claims that lack a declared baseline year.", GOLD_ACCENT, 6.8, 1.8),
        ("Greenwashing Risk Score", "Objective 0-100 risk score (Low, Medium, High Risk).", NEON_CYAN, 0.8, 4.4),
        ("Auditor Recommendations", "Generates concrete corrective action steps for compliance officers.", EMERALD_GREEN, 6.8, 4.4)
    ]

    for title, desc, color, left, top in risk_items:
        card = add_card(s10, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s10, card)

    set_speaker_notes(
        s10,
        say="Our Risk Agent exposes greenwashing by flagging vague terms, missing baseline years, and unbacked claims, producing an actionable risk score.",
        click="Click 4 times — one per risk component.",
        pause="Pause on Vague Term Detection to give examples like '100% eco-friendly'.",
        judges="Judges will recognize the anti-greenwashing core of our platform."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 11: Dashboard Overview
    # ───────────────────────────────────────────────────────────────────────────
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, 11, "Product Experience", "Centralized Overview Dashboard UI")

    dash_widgets = [
        ("Top 4 KPI Scorecards", "Portfolio EcoScore, Carbon Footprint, Active Labels, SKUs.", EMERALD_GREEN, 0.8, 1.8),
        ("AI Agent Live Status", "Real-time task monitor for 4 active backend agents.", NEON_CYAN, 6.8, 1.8),
        ("Trust & Carbon Widgets", "Interactive Trust Score ring & 12-month carbon trend.", PURPLE_ACCENT, 0.8, 4.4),
        ("Upload & History Card", "Instant drag-and-drop PDF dropzone + upload history.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in dash_widgets:
        card = add_card(s11, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s11, card)

    set_speaker_notes(
        s11,
        say="The Overview Dashboard provides executive visibility — displaying EcoScore, Carbon Footprint, active AI agents, and single-click PDF upload.",
        click="Click 4 times to reveal each UI widget zone.",
        pause="Pause after revealing all 4 zones to showcase design polish.",
        judges="Judges will notice the clean glassmorphic UI layout."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 12: Products Module
    # ───────────────────────────────────────────────────────────────────────────
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, 12, "Product Intelligence", "Automated Product & SKU Extraction")

    prod_items = [
        ("Automated Extraction", "Extracts product names & SKUs directly from PDF disclosures.", EMERALD_GREEN, 0.8, 1.8),
        ("Category Classification", "Groups items into Electronics, Packaging, Agriculture & Materials.", NEON_CYAN, 6.8, 1.8),
        ("Carbon Impact Per Unit", "Assigns estimated product carbon intensity footprint.", PURPLE_ACCENT, 0.8, 4.4),
        ("Verification Status", "Displays Verified vs Pending review status per product.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in prod_items:
        card = add_card(s12, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s12, card)

    set_speaker_notes(
        s12,
        say="In the Products Module, EcoLabel X automatically extracts individual product SKUs, assigns carbon footprint metrics, and tracks verification status.",
        click="Click 4 times to reveal product extraction steps.",
        pause="Pause on Carbon Impact to emphasize item-level ESG granular tracking.",
        judges="Judges will see how ESG data connects to physical product SKUs."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 13: Claim Detector
    # ───────────────────────────────────────────────────────────────────────────
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, 13, "Claim Detector", "NLP & Keyword Claim Extraction Pipeline")

    cd_items = [
        ("Exact Quote Parsing", "Extracts exact environmental text snippets from disclosures.", EMERALD_GREEN, 0.8, 1.8),
        ("Page Number Mapping", "Pinpoints exact PDF page location for auditor inspection.", NEON_CYAN, 6.8, 1.8),
        ("Keyword Weighting", "Identifies matched terms like 'RE100', 'net zero', 'recycled'.", PURPLE_ACCENT, 0.8, 4.4),
        ("Confidence Scoring", "Assigns AI confidence score (0-100%) to every extracted claim.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in cd_items:
        card = add_card(s13, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s13, card)

    set_speaker_notes(
        s13,
        say="Our Claim Detector extracts exact text quotes from disclosures and pinpoints their exact page number for instant audit verification.",
        click="Click 4 times — one per claim extraction stage.",
        pause="Pause on Page Number Mapping to highlight audit traceability.",
        judges="Judges will note exact page traceability for corporate compliance."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 14: Evidence Verification
    # ───────────────────────────────────────────────────────────────────────────
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, 14, "Evidence Verifier", "Automated ISO & Standards Verification")

    ev_items = [
        ("Evidence Search", "Searches quantitative data within the report to back claims.", EMERALD_GREEN, 0.8, 1.8),
        ("ISO Rule Checking", "Validates claims against ISO 14021 environmental self-declarations.", NEON_CYAN, 6.8, 1.8),
        ("Rationale Generation", "Generates clear AI explanation for why a claim is verified or flagged.", PURPLE_ACCENT, 0.8, 4.4),
        ("Verification Badge", "Outputs Verified, Under Review, or Unverified status badges.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in ev_items:
        card = add_card(s14, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s14, card)

    set_speaker_notes(
        s14,
        say="The Evidence Verifier validates claims against ISO standards and provides clear written rationales explaining why a claim passed or failed.",
        click="Click 4 times to reveal evidence verification flow.",
        pause="Pause on Rationale Generation to explain explainable AI logic.",
        judges="Judges will appreciate transparent, explainable AI decisions."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 15: Greenwashing Analyzer
    # ───────────────────────────────────────────────────────────────────────────
    s15 = prs.slides.add_slide(blank_layout)
    add_header(s15, 15, "Risk Analyzer", "Greenwashing Risk Matrix & Recommendations")

    gw_items = [
        ("Risk Level Classification", "Categorizes overall risk into Low, Medium, or High Risk.", RED_ACCENT, 0.8, 1.8),
        ("Risk Driver Breakdown", "Lists specific reasons (e.g. unverified packaging claim).", GOLD_ACCENT, 6.8, 1.8),
        ("Missing Evidence List", "Flags missing third-party certificates or audit documents.", NEON_CYAN, 0.8, 4.4),
        ("Actionable Guidance", "Provides prioritized recommendation steps for risk mitigation.", EMERALD_GREEN, 6.8, 4.4)
    ]

    for title, desc, color, left, top in gw_items:
        card = add_card(s15, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s15, card)

    set_speaker_notes(
        s15,
        say="The Greenwashing Analyzer categorizes overall document risk, pinpoints missing evidence, and provides concrete mitigation steps.",
        click="Click 4 times to reveal risk analyzer components.",
        pause="Pause on Missing Evidence List to show how missing proof is detected.",
        judges="Judges will see actionable risk mitigation for enterprise leaders."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 16: Carbon Dashboard
    # ───────────────────────────────────────────────────────────────────────────
    s16 = prs.slides.add_slide(blank_layout)
    add_header(s16, 16, "Carbon Analytics", "Scope 1, 2 & 3 Footprint Visualization")

    car_items = [
        ("Scope 1 Direct", "Onsite manufacturing & facility emissions (kt CO2e).", EMERALD_GREEN, 0.8, 1.8),
        ("Scope 2 Energy", "Purchased electricity & heating emissions (kt CO2e).", NEON_CYAN, 6.8, 1.8),
        ("Scope 3 Value Chain", "Raw material, freight & distribution emissions (kt CO2e).", PURPLE_ACCENT, 0.8, 4.4),
        ("Intensity Reductions", "Visual progress bars tracking YoY carbon reduction targets.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in car_items:
        card = add_card(s16, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s16, card)

    set_speaker_notes(
        s16,
        say="On the Carbon Dashboard, users explore interactive Scope 1, 2, and 3 emissions breakdowns and monitor carbon reduction progress.",
        click="Click 4 times — reveal each emission scope component.",
        pause="Pause after all 4 cards appear to highlight full GHG scope coverage.",
        judges="Judges will notice comprehensive carbon accounting alignment."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 17: Analytics Dashboard
    # ───────────────────────────────────────────────────────────────────────────
    s17 = prs.slides.add_slide(blank_layout)
    add_header(s17, 17, "Compliance Analytics", "Category Distribution & Audit Index")

    ana_items = [
        ("Category Distribution", "Visual breakdown of claims across Energy, Carbon, Packaging, and Supply Chain.", EMERALD_GREEN, 0.8, 1.8),
        ("Audit Compliance Index", "Dynamic score evaluating EU ESPR & FTC Green Guide readiness.", NEON_CYAN, 6.8, 1.8),
        ("Verification Rate", "Displays percentage of total claims backed by quantitative proof.", PURPLE_ACCENT, 0.8, 4.4),
        ("Audit Readiness Rating", "Categorizes overall audit readiness from High to Action Needed.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in ana_items:
        card = add_card(s17, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s17, card)

    set_speaker_notes(
        s17,
        say="Our Analytics Module calculates claim category distributions and computes an Audit Compliance Index score for regulatory readiness.",
        click="Click 4 times to build the analytics view.",
        pause="Pause on Audit Compliance Index to explain compliance scoring logic.",
        judges="Judges will observe high-level executive analytics paired with granular metrics."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 18: Eco Labels Registry
    # ───────────────────────────────────────────────────────────────────────────
    s18 = prs.slides.add_slide(blank_layout)
    add_header(s18, 18, "Certification Registry", "Automated Eco Label Identification")

    labels_list = [
        ("ISO 14064 GHG Verified", "TÜV SÜD / WRI Emissions Assurance", EMERALD_GREEN, 0.8, 1.8),
        ("FSC Certified Packaging", "Forest Stewardship Council Chain of Custody", NEON_CYAN, 6.8, 1.8),
        ("RE100 Renewable Power", "Climate Group 100% Clean Energy Mark", PURPLE_ACCENT, 0.8, 4.4),
        ("EU Organic & Fairtrade", "European Commission & Fairtrade International", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in labels_list:
        card = add_card(s18, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s18, card)

    set_speaker_notes(
        s18,
        say="The Eco Labels Registry automatically detects third-party certifications (ISO 14064, FSC, RE100, EU Organic) from disclosures, verifying validity and page excerpts.",
        click="Click 4 times — one per eco label category.",
        pause="Pause on ISO 14064 to explain GHG verification standard matching.",
        judges="Judges will see how certification detection eliminates manual registry searching."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 19: Executive Audit Report Exporter
    # ───────────────────────────────────────────────────────────────────────────
    s19 = prs.slides.add_slide(blank_layout)
    add_header(s19, 19, "Executive Reporting", "Publication-Ready A4 PDF Report Exporter")

    rep_items = [
        ("Official Letterhead", "EcoLabel X header with ISO 14021 compliance seal & metadata.", EMERALD_GREEN, 0.8, 1.8),
        ("KPI Scorecards", "Overall Trust Score, Risk Score, Claims Ratio scorecard grid.", NEON_CYAN, 6.8, 1.8),
        ("Claims Directory Table", "Complete table of verified & flagged claims with PDF page numbers.", PURPLE_ACCENT, 0.8, 4.4),
        ("Auditor Sign-off Block", "Prioritized recommendations & formal auditor signature section.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in rep_items:
        card = add_card(s19, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s19, card)

    set_speaker_notes(
        s19,
        say="With one click on 'Download PDF', EcoLabel X exports a publication-ready A4 executive audit document complete with letterhead, scorecards, and auditor sign-off.",
        click="Click 4 times to reveal PDF report sections.",
        pause="Pause after all 4 cards appear to emphasize PDF export polish.",
        judges="Judges will see an end-to-end usable product that delivers a tangible PDF artifact."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 20: Technology Stack
    # ───────────────────────────────────────────────────────────────────────────
    s20 = prs.slides.add_slide(blank_layout)
    add_header(s20, 20, "Engineering Core", "Modern Full-Stack Technology Suite")

    tech_items = [
        ("Next.js 14 App Router", "React, TypeScript, Glassmorphic UI framework", EMERALD_GREEN, 0.8, 1.8),
        ("FastAPI Backend", "Async Python 3.10 server with Pydantic v2 schemas", NEON_CYAN, 6.8, 1.8),
        ("Google Gemini 2.5 Flash", "Generative AI LLM reasoning & claim analysis", PURPLE_ACCENT, 0.8, 4.4),
        ("PyPDF2 & PDF Exporter", "Raw PDF extraction & A4 publication exporter", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in tech_items:
        card = add_card(s20, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s20, card)

    set_speaker_notes(
        s20,
        say="Our technology stack leverages Next.js 14, FastAPI, Python, PyPDF2, and Google Gemini 2.5 Flash API for high-speed, enterprise-grade AI execution.",
        click="Click 4 times — one per technology pillar.",
        pause="Pause on Gemini 2.5 Flash to highlight LLM integration.",
        judges="Judges will recognize production-ready, modern open-source technologies."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 21: Key Innovations
    # ───────────────────────────────────────────────────────────────────────────
    s21 = prs.slides.add_slide(blank_layout)
    add_header(s21, 21, "Competitive Edge", "Why EcoLabel X Is Unique")

    inno_items = [
        ("1.5-Second Execution", "Replaces days of manual PDF auditing with instant analysis.", EMERALD_GREEN, 0.8, 1.8),
        ("4 Parallel AI Agents", "Multi-agent architecture evaluating distinct compliance vectors.", NEON_CYAN, 6.8, 1.8),
        ("Evidence-Backed Scoring", "Every score is linked to exact text quotes and PDF page numbers.", PURPLE_ACCENT, 0.8, 4.4),
        ("One-Click Executive PDF", "Generates publication-ready A4 executive audit reports.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in inno_items:
        card = add_card(s21, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s21, card)

    set_speaker_notes(
        s21,
        say="What makes EcoLabel X unique is sub-2-second speed, 4-agent parallel evaluation, 100% page-traceable evidence, and instant executive PDF generation.",
        click="Click 4 times to reveal key innovations.",
        pause="Pause on Evidence-Backed Scoring to reinforce trust.",
        judges="Judges will see clear differentiation from simple chatbots or basic scanners."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 22: Future Scope & Roadmap
    # ───────────────────────────────────────────────────────────────────────────
    s22 = prs.slides.add_slide(blank_layout)
    add_header(s22, 22, "Product Vision", "Future Roadmap & Strategic Scope")

    road_items = [
        ("OCR for Scanned PDFs", "Integrating Tesseract/Vision API for scanned legacy documents.", EMERALD_GREEN, 0.8, 1.8),
        ("Multi-Language Support", "Auditing ESG reports in German, French, Spanish & Mandarin.", NEON_CYAN, 6.8, 1.8),
        ("Live Regulatory Feeds", "Real-time updates as EU ESPR & FTC regulations evolve.", PURPLE_ACCENT, 0.8, 4.4),
        ("Enterprise Peer Benchmarking", "Comparing ESG performance across industry competitors.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in road_items:
        card = add_card(s22, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s22, card)

    set_speaker_notes(
        s22,
        say="Our future roadmap includes OCR for scanned PDF disclosures, multi-language support, live regulatory feeds, and industry peer benchmarking.",
        click="Click 4 times — reveal each roadmap item.",
        pause="Pause on Multi-Language Support to mention global ESG expansion.",
        judges="Judges will appreciate a visionary yet practical growth strategy."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 23: Live Demo Flow
    # ───────────────────────────────────────────────────────────────────────────
    s23 = prs.slides.add_slide(blank_layout)
    add_header(s23, 23, "Live Demonstration", "Step-by-Step Live Walkthrough Flow")

    demo_steps = [
        ("Step 1: Select & Upload", "Drag PDF report into dropzone", EMERALD_GREEN, 0.8, 1.8),
        ("Step 2: Instant Analysis", "Backend parses PDF in <1.5s", NEON_CYAN, 6.8, 1.8),
        ("Step 3: Dashboard Insights", "Explore EcoScore & Carbon Cards", PURPLE_ACCENT, 0.8, 4.4),
        ("Step 4: Export Audit PDF", "Download publication-ready report", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in demo_steps:
        card = add_card(s23, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s23, card)

    set_speaker_notes(
        s23,
        say="Let us demonstrate EcoLabel X live. Watch as we upload a PDF, run 4 AI agents in <1.5s, explore dashboard metrics, and download the executive audit report.",
        click="Click 4 times as you transition to the live application demo.",
        pause="Pause on Step 4 before switching to http://localhost:3000.",
        judges="Judges will experience an impressive, seamless live product demonstration."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 24: Business & Environmental Impact
    # ───────────────────────────────────────────────────────────────────────────
    s24 = prs.slides.add_slide(blank_layout)
    add_header(s24, 24, "Value Created", "Quantifiable Business & ESG Impact")

    impact_items = [
        ("95% Audit Time Reduction", "Reduces ESG report analysis from weeks to 1.5 seconds.", EMERALD_GREEN, 0.8, 1.8),
        ("100% Audit Transparency", "Every score is linked to exact quotes and PDF page numbers.", NEON_CYAN, 6.8, 1.8),
        ("Zero Regulatory Penalties", "Prevents costly EU ESPR & FTC Green Guide non-compliance fines.", PURPLE_ACCENT, 0.8, 4.4),
        ("Real Carbon Accountability", "Accelerates genuine Scope 1-3 reduction across supply chains.", GOLD_ACCENT, 6.8, 4.4)
    ]

    for title, desc, color, left, top in impact_items:
        card = add_card(s24, left, top, 5.6, 2.2, title, desc, border=color)
        add_animation(s24, card)

    set_speaker_notes(
        s24,
        say="EcoLabel X delivers measurable ROI: 95% audit time savings, 100% page-level transparency, and zero regulatory greenwashing penalties.",
        click="Click 4 times to reveal each impact metric.",
        pause="Pause on 95% Audit Time Reduction to leave a lasting quantitative metric.",
        judges="Judges will see tangible business ROI and real environmental value."
    )

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 25: Conclusion & Thank You
    # ───────────────────────────────────────────────────────────────────────────
    s25 = prs.slides.add_slide(blank_layout)
    set_background(s25)

    c25 = add_card(s25, 1.5, 1.2, 10.33, 5.1, bg=RGBColor(15, 23, 42), border=EMERALD_GREEN)
    tf25 = c25.text_frame
    tf25.word_wrap = True
    tf25.margin_top = Inches(0.8)

    p25 = tf25.paragraphs[0]
    p25.text = "Thank You!"
    p25.font.size = Pt(54)
    p25.font.bold = True
    p25.font.color.rgb = EMERALD_GREEN
    p25.alignment = PP_ALIGN.CENTER

    p25_sub = tf25.add_paragraph()
    p25_sub.text = "EcoLabel X — Transforming Sustainability Intelligence"
    p25_sub.font.size = Pt(20)
    p25_sub.font.bold = True
    p25_sub.font.color.rgb = NEON_CYAN
    p25_sub.alignment = PP_ALIGN.CENTER
    p25_sub.space_before = Pt(16)

    p25_urls = tf25.add_paragraph()
    p25_urls.text = "🌐 Live Demo: http://localhost:3000    •    📦 GitHub: github.com/mohamedidhris777/eco-label"
    p25_urls.font.size = Pt(13)
    p25_urls.font.color.rgb = TEXT_MUTED
    p25_urls.alignment = PP_ALIGN.CENTER
    p25_urls.space_before = Pt(28)

    p25_qa = tf25.add_paragraph()
    p25_qa.text = "WE ARE NOW OPEN FOR QUESTIONS"
    p25_qa.font.size = Pt(12)
    p25_qa.font.bold = True
    p25_qa.font.color.rgb = GOLD_ACCENT
    p25_qa.alignment = PP_ALIGN.CENTER
    p25_qa.space_before = Pt(36)

    set_speaker_notes(
        s25,
        say="Thank you judges for your time! We are now open for your questions.",
        click="No further clicks needed. Leave slide open for Q&A.",
        pause="Pause and look at the judges confidently.",
        judges="Judges can view live app URL and GitHub repository links while asking questions."
    )

    out_file1 = r"c:\Users\moham\.gemini\antigravity-ide\scratch\idhris-ai\ecolabel-x\EcoLabel_X_Hackathon_Presentation.pptx"
    out_file2 = r"C:\Users\moham\.gemini\antigravity-ide\brain\bef9a80d-3767-499d-99bf-28b0c7a9f8a9\EcoLabel_X_Hackathon_Presentation.pptx"

    prs.save(out_file1)
    prs.save(out_file2)
    print(f"SUCCESS: Saved PowerPoint deck to:\n1. {out_file1}\n2. {out_file2}")

if __name__ == "__main__":
    create_deck()
